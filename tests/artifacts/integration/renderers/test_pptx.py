"""Integration tests for the PPTX renderer.

Marker: ``integration`` (no real LLM needed — hand-constructed schemas).

These tests:
1. Render a 5-slide deck (cover + 4 content slides) and assert slide count.
2. Assert slide 0 is a title-only cover with NO bullets on the body placeholder.
3. Assert slides 1-4 each have bullets on the body placeholder.
4. Verify the themed template is used (presentation dimensions match 13.33×7.5).
5. Test embed_chart: a dummy chart PNG is embedded on a dedicated slide.
6. Test PitchDeckSchema (tagline key) round-trip.
"""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from open_notebook.artifacts.generators.slide_deck import SlideDeckSchema, SlideSchema
from open_notebook.artifacts.generators.pitch_deck import PitchDeckSchema, PitchSlideSchema
from open_notebook.artifacts.renderers.pptx_renderer import render_deck

pytestmark = pytest.mark.integration


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def five_slide_deck() -> SlideDeckSchema:
    """A SlideDeckSchema with 4 content slides (cover is auto-generated)."""
    return SlideDeckSchema(
        title="SOTA Artifact Pipeline",
        subtitle="Phase 2 Stream D",
        slides=[
            SlideSchema(
                title="Introduction",
                bullets=["Open source", "Privacy-first", "Multi-provider"],
                notes="Welcome notes here",
            ),
            SlideSchema(
                title="Architecture",
                bullets=["FastAPI backend", "SurrealDB", "LangGraph"],
                notes="Cover three-tier design",
            ),
            SlideSchema(
                title="Generators",
                bullets=["13 artifact types", "Structured output", "Pydantic schemas"],
            ),
            SlideSchema(
                title="Conclusion",
                bullets=["Real LLM calls only", "No heuristics", "Production quality"],
                notes="Call to action",
            ),
        ],
    )


@pytest.fixture
def pitch_deck() -> PitchDeckSchema:
    """A PitchDeckSchema with canonical VC slides."""
    return PitchDeckSchema(
        title="OpenNotebook AI",
        tagline="The privacy-first NotebookLM alternative",
        slides=[
            PitchSlideSchema(title="Problem", bullets=["Data silos", "Vendor lock-in"]),
            PitchSlideSchema(title="Solution", bullets=["Open source", "Self-hosted"]),
            PitchSlideSchema(title="Market", bullets=["$10B TAM", "Growing SMB segment"]),
            PitchSlideSchema(title="Ask", bullets=["$2M seed", "12-month runway"]),
        ],
    )


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestPptxSlideCount:
    def test_slide_count_equals_schema_slides_plus_cover(
        self, five_slide_deck: SlideDeckSchema, tmp_path: Path
    ) -> None:
        """Rendered deck must have exactly len(slides) + 1 slides (cover + content)."""
        output = tmp_path / "deck.pptx"
        render_deck(five_slide_deck, output)

        prs = Presentation(str(output))
        # 4 content slides + 1 cover = 5
        assert len(prs.slides) == 5, (
            f"Expected 5 slides (1 cover + 4 content), got {len(prs.slides)}"
        )


class TestPptxCoverSlide:
    def test_slide_0_is_title_only(
        self, five_slide_deck: SlideDeckSchema, tmp_path: Path
    ) -> None:
        """Slide 0 (cover) must have title and subtitle placeholder only.

        The body placeholder must NOT contain any of the first content slide's
        bullets (fixes the Phase 0 bullet-loss / slide-0 hijack bug).
        """
        output = tmp_path / "cover.pptx"
        render_deck(five_slide_deck, output)

        prs = Presentation(str(output))
        cover = prs.slides[0]

        # Title placeholder must match deck title
        assert cover.shapes.title is not None
        assert five_slide_deck.title in cover.shapes.title.text

        # Verify that none of the first content slide's bullets appear on the cover
        first_content_slide_bullets = five_slide_deck.slides[0].bullets
        cover_text_all = " ".join(
            shape.text_frame.text
            for shape in cover.shapes
            if shape.has_text_frame
        )
        for bullet in first_content_slide_bullets:
            assert bullet not in cover_text_all, (
                f"Bullet {bullet!r} from slide[0] appeared on the cover slide. "
                "This is the slide-0 bullet-loss bug — fix render_deck()."
            )


class TestPptxContentSlides:
    def test_content_slides_have_bullets(
        self, five_slide_deck: SlideDeckSchema, tmp_path: Path
    ) -> None:
        """Each content slide (indices 1-N) must have at least one bullet in body."""
        output = tmp_path / "content.pptx"
        render_deck(five_slide_deck, output)

        prs = Presentation(str(output))
        # Content slides are indices 1..N-1
        for idx in range(1, len(prs.slides)):
            slide = prs.slides[idx]
            schema_slide = five_slide_deck.slides[idx - 1]  # offset by cover

            # Find the body placeholder (idx != 0)
            body_placeholder = next(
                (ph for ph in slide.placeholders
                 if ph.placeholder_format.idx != 0),
                None,
            )
            assert body_placeholder is not None, (
                f"Content slide {idx} has no body placeholder."
            )

            body_text = body_placeholder.text_frame.text
            # At least one expected bullet must be in the body
            found = any(
                bullet in body_text
                for bullet in schema_slide.bullets
            )
            assert found, (
                f"Slide {idx} body placeholder is missing expected bullets. "
                f"Expected one of {schema_slide.bullets!r}, got: {body_text!r}"
            )

    def test_speaker_notes_written(
        self, five_slide_deck: SlideDeckSchema, tmp_path: Path
    ) -> None:
        """Slides with notes in schema must have notes in the pptx notes_text_frame."""
        output = tmp_path / "notes.pptx"
        render_deck(five_slide_deck, output)

        prs = Presentation(str(output))
        for idx, schema_slide in enumerate(five_slide_deck.slides, start=1):
            if schema_slide.notes:
                notes_tf = prs.slides[idx].notes_slide.notes_text_frame
                assert schema_slide.notes in notes_tf.text, (
                    f"Notes for slide {idx} missing in pptx notes_text_frame."
                )


class TestPptxPresentationDimensions:
    def test_dimensions_are_widescreen(
        self, five_slide_deck: SlideDeckSchema, tmp_path: Path
    ) -> None:
        """Presentation must use widescreen 13.33 × 7.5 inch slide size."""
        output = tmp_path / "dims.pptx"
        render_deck(five_slide_deck, output)

        prs = Presentation(str(output))
        from pptx.util import Inches
        width_inches = prs.slide_width / 914400  # EMUs per inch
        height_inches = prs.slide_height / 914400

        assert abs(width_inches - 13.333) < 0.1, (
            f"Expected width ~13.333 inches, got {width_inches:.3f}"
        )
        assert abs(height_inches - 7.5) < 0.1, (
            f"Expected height ~7.5 inches, got {height_inches:.3f}"
        )


class TestPptxPitchDeck:
    def test_pitch_deck_tagline_on_cover(
        self, pitch_deck: PitchDeckSchema, tmp_path: Path
    ) -> None:
        """Pitch deck cover must contain company name and tagline."""
        output = tmp_path / "pitch.pptx"
        render_deck(pitch_deck, output, subtitle_key="tagline")

        prs = Presentation(str(output))
        cover = prs.slides[0]
        cover_text = " ".join(
            shape.text_frame.text
            for shape in cover.shapes
            if shape.has_text_frame
        )
        assert pitch_deck.title in cover_text
        assert pitch_deck.tagline in cover_text

    def test_pitch_deck_slide_count(
        self, pitch_deck: PitchDeckSchema, tmp_path: Path
    ) -> None:
        output = tmp_path / "pitch2.pptx"
        render_deck(pitch_deck, output, subtitle_key="tagline")

        prs = Presentation(str(output))
        assert len(prs.slides) == len(pitch_deck.slides) + 1  # cover + content


class TestPptxChartEmbedding:
    def test_embed_chart_adds_extra_slide(
        self, five_slide_deck: SlideDeckSchema, tmp_path: Path
    ) -> None:
        """When embed_chart is provided, an extra image slide is appended."""
        # Create a minimal PNG via Pillow
        from PIL import Image
        chart_path = tmp_path / "chart.png"
        img = Image.new("RGB", (800, 600), color=(35, 102, 194))
        img.save(str(chart_path))

        output = tmp_path / "with_chart.pptx"
        render_deck(five_slide_deck, output, embed_chart=chart_path)

        prs = Presentation(str(output))
        # cover + 4 content + 1 chart = 6
        assert len(prs.slides) == 6, (
            f"Expected 6 slides (cover + 4 content + 1 chart), got {len(prs.slides)}"
        )
