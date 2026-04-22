"""python-pptx renderer for slide and pitch decks.

SOTA rewrite for Phase 2 Stream D.

Key improvements over the Phase 0 version:
- Loads a themed master template from ``assets/artifacts/themes/default.pptx``
  so every deck shares a consistent professional colour palette and typography.
- **Fixes the slide-0 bullet-loss bug**: slide 0 is *always* a proper cover slide
  containing only the deck ``title`` + ``subtitle`` / ``tagline``.  The first
  :class:`~open_notebook.artifacts.generators.slide_deck.SlideSchema` in the
  ``slides`` list is rendered as slide 1 (a content slide), not as the cover.
- Subsequent slides render ``title + bullets + speaker notes`` via the
  "Title and Content" layout (index 1).
- Speaker notes are written to the ``notes_slide.notes_text_frame``.
- Supports embedding a chart image (PNG path) at a designated image placeholder
  on the last content slide via :func:`embed_chart_image`.
- Two public entry points:
  - :func:`render_deck` — accepts a ``SlideDeckSchema | PitchDeckSchema``
    *model_dump()* dict or a Pydantic model instance and writes a ``.pptx``.
  - :func:`embed_image_on_slide` — lower-level helper called by Stream G's
    image_gen pipeline to stamp a PNG image onto an existing slide.

Layout index reference (from default.pptx):
  0 = Title Slide          (cover: title + subtitle)
  1 = Title and Content    (body slides: title + bulleted body)
  5 = Title Only           (image/chart slide: title, no body placeholder)
  6 = Blank                (used if no layout fits)
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme constants (mirror the colours defined in default.pptx theme XML)
# ---------------------------------------------------------------------------

_BRAND_DARK = RGBColor(0x1E, 0x29, 0x3B)    # navy title text
_BRAND_BLUE = RGBColor(0x23, 0x66, 0xC2)    # accent / headings
_BRAND_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)   # near-white for dark backgrounds
_BRAND_GRAY = RGBColor(0x64, 0x74, 0x8B)    # subtitle / notes

# ---------------------------------------------------------------------------
# Template path
# ---------------------------------------------------------------------------

_THEME_PATH: Path = (
    Path(__file__).parent.parent.parent.parent  # repo root
    / "assets" / "artifacts" / "themes" / "default.pptx"
)


def _load_presentation() -> Presentation:
    """Load the branded theme template, falling back to blank if absent."""
    if _THEME_PATH.exists():
        prs = Presentation(str(_THEME_PATH))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_cover_slide(prs: Presentation, title: str, subtitle: Optional[str]) -> None:
    """Add layout-0 *Title Slide* (cover only — no bullet body)."""
    layout = prs.slide_layouts[0]  # "Title Slide"
    slide = prs.slides.add_slide(layout)

    # Title placeholder (idx 0)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title or "Presentation"
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = _BRAND_DARK

    # Subtitle placeholder (idx 1) — subtitle text only, never bullets
    if len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle or ""
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = _BRAND_GRAY


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    notes: Optional[str] = None,
) -> None:
    """Add layout-1 *Title and Content* slide with bullets and optional notes."""
    layout = prs.slide_layouts[1]  # "Title and Content"
    slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title or "Slide"
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = _BRAND_BLUE

    # Body / bullets
    body_ph = next(
        (ph for ph in slide.placeholders if ph.placeholder_format.idx != 0),
        None,
    )
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        bullet_list = bullets if bullets else ["—"]
        for idx, bullet in enumerate(bullet_list):
            if idx == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = str(bullet)
            para.level = 0
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = _BRAND_DARK

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


def _add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    notes: Optional[str] = None,
) -> None:
    """Add a *Title Only* slide (layout 5) with a full-body chart/image."""
    # Use "Title Only" layout if available, else Blank
    try:
        layout = prs.slide_layouts[5]  # "Title Only"
    except IndexError:
        layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title or "Figure"

    # Place image covering the content area (below title)
    img_top = Inches(1.5)
    img_left = Inches(0.5)
    img_width = Inches(12.0)
    img_height = Inches(5.5)
    slide.shapes.add_picture(
        str(image_path),
        img_left,
        img_top,
        width=img_width,
        height=img_height,
    )

    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def embed_image_on_slide(slide: Any, image_path: Path, *, left: float = 0.5,
                         top: float = 1.5, width: float = 12.0,
                         height: float = 5.5) -> None:
    """Stamp *image_path* onto an existing slide object.

    Lower-level helper used by Stream G's video/image-gen pipeline.  Caller
    is responsible for obtaining the slide reference.

    Args:
        slide:      A ``pptx.slide.Slide`` instance.
        image_path: Path to a PNG/JPEG file to embed.
        left:       Left offset in inches (default 0.5).
        top:        Top offset in inches (default 1.5, clears title bar).
        width:      Image width in inches (default 12.0).
        height:     Image height in inches (default 5.5).
    """
    slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )


def render_deck(
    data: Union[Dict[str, Any], Any],
    path: Path,
    *,
    subtitle_key: str = "subtitle",
    embed_chart: Optional[Path] = None,
) -> Path:
    """Render a slide or pitch deck to *path* as a ``.pptx`` file.

    Accepts either a plain ``dict`` (from ``schema.model_dump()``) or a
    Pydantic model instance (``SlideDeckSchema`` or ``PitchDeckSchema``).

    The **slide-0 bullet-loss bug** (Phase 0 audit finding) is fixed here:
    - Slide 0 is *always* the cover (title + subtitle only).
    - Every entry in ``data["slides"]`` is rendered as a content slide
      starting at position 1 — including the first slide, which previously
      had its bullets discarded and was wrongly treated as the cover.

    Args:
        data:         ``model_dump()`` dict or Pydantic schema instance.
        path:         Output file path (parent dirs created automatically).
        subtitle_key: Dict key for the subtitle / tagline field
                      (``"subtitle"`` for SlideDeck, ``"tagline"`` for Pitch).
        embed_chart:  Optional path to a PNG to embed on the last slide as a
                      "Title Only" chart slide appended after the content slides.

    Returns:
        The resolved output path.
    """
    # Accept Pydantic model instances as well as plain dicts
    if not isinstance(data, dict):
        data = data.model_dump()

    prs = _load_presentation()

    deck_title: str = data.get("title", "Presentation") or "Presentation"
    deck_subtitle: Optional[str] = data.get(subtitle_key) or data.get("tagline")
    slides: List[Dict[str, Any]] = data.get("slides") or []

    # --- Cover slide (slide 0) -------------------------------------------
    _add_cover_slide(prs, deck_title, deck_subtitle)

    # --- Content slides (slides 1…N) — ALL slides from schema --------------
    # BUG FIX: Phase 0 renderer used slides[1:] here, discarding the first
    # slide's bullets.  We now render ALL slides as content slides, after a
    # dedicated cover slide.
    for slide in slides:
        _add_content_slide(
            prs,
            title=slide.get("title", ""),
            bullets=slide.get("bullets") or [],
            notes=slide.get("notes"),
        )

    # --- Optional chart/image slide ----------------------------------------
    if embed_chart is not None and Path(embed_chart).exists():
        _add_image_slide(prs, "Figure", Path(embed_chart))

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
